"""
attachment_tool.py
-------------------
Recebe qualquer ficheiro anexado no chat (documentos, folhas de cálculo,
apresentações, código-fonte, ficheiros comprimidos, áudio, vídeo, etc.) e
tenta extrair texto/conteúdo legível para dar como contexto ao modelo de IA,
para que este consiga responder a perguntas sobre o ficheiro.

Imagens continuam a ser tratadas à parte pelo ocr_tool.py (criação do PDF
pesquisável), este módulo cobre todos os outros tipos.

Cada tipo de ficheiro tem o seu próprio extrator "best effort": se a
extração falhar ou o tipo não for suportado, devolve-se None e o ficheiro
é apenas referenciado pelo nome, para que o modelo saiba que foi anexado
mesmo sem conseguir "ler" o conteúdo (ex.: áudio/vídeo, binários).
"""

import csv
import io
import os
import zipfile

MAX_CHARS_PER_FILE = 4000

TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".tsv", ".json", ".xml", ".yaml", ".yml",
    ".log", ".ini", ".cfg", ".conf",
    # código-fonte
    ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".cpp", ".h",
    ".hpp", ".cs", ".go", ".rs", ".rb", ".php", ".sh", ".bat", ".ps1",
    ".sql", ".html", ".css", ".scss", ".vue", ".swift", ".kt",
}


def _truncate(texto: str) -> str:
    texto = texto.strip()
    if len(texto) > MAX_CHARS_PER_FILE:
        return texto[:MAX_CHARS_PER_FILE] + "\n[...conteúdo truncado...]"
    return texto


def _extract_plain_text(raw: bytes) -> str:
    return raw.decode("utf-8", errors="ignore")


def _extract_pdf(raw: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(raw))
    partes = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(partes)


def _extract_docx(raw: bytes) -> str:
    import docx
    documento = docx.Document(io.BytesIO(raw))
    partes = [p.text for p in documento.paragraphs if p.text]
    for tabela in documento.tables:
        for linha in tabela.rows:
            partes.append(" | ".join(c.text for c in linha.cells))
    return "\n".join(partes)


def _extract_pptx(raw: bytes) -> str:
    from pptx import Presentation
    apresentacao = Presentation(io.BytesIO(raw))
    partes = []
    for i, slide in enumerate(apresentacao.slides, start=1):
        textos_slide = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                textos_slide.append(shape.text)
        if textos_slide:
            partes.append(f"[Slide {i}]\n" + "\n".join(textos_slide))
    return "\n\n".join(partes)


def _extract_xlsx(raw: bytes) -> str:
    import openpyxl
    livro = openpyxl.load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
    partes = []
    for nome_folha in livro.sheetnames:
        folha = livro[nome_folha]
        linhas_texto = []
        for i, linha in enumerate(folha.iter_rows(values_only=True)):
            if i >= 200:  # limite de segurança por folha
                linhas_texto.append("[...mais linhas omitidas...]")
                break
            valores = ["" if v is None else str(v) for v in linha]
            if any(valores):
                linhas_texto.append(" | ".join(valores))
        if linhas_texto:
            partes.append(f"[Folha: {nome_folha}]\n" + "\n".join(linhas_texto))
    return "\n\n".join(partes)


def _extract_zip(raw: bytes) -> str:
    with zipfile.ZipFile(io.BytesIO(raw)) as z:
        nomes = z.namelist()
    listagem = "\n".join(f"- {n}" for n in nomes[:200])
    if len(nomes) > 200:
        listagem += f"\n[...e mais {len(nomes) - 200} ficheiro(s)...]"
    return f"Arquivo comprimido com {len(nomes)} ficheiro(s):\n{listagem}"


def extract_attachment_text(filename: str, mime: str, raw: bytes) -> str | None:
    """
    Tenta extrair conteúdo legível de um ficheiro anexado (não-imagem).

    Devolve o texto extraído (truncado a MAX_CHARS_PER_FILE) ou None se o
    tipo de ficheiro não for suportado / a extração falhar.
    """
    ext = os.path.splitext(filename or "")[1].lower()
    mime = (mime or "").lower()

    try:
        if ext == ".pdf" or mime == "application/pdf":
            texto = _extract_pdf(raw)
        elif ext == ".docx":
            texto = _extract_docx(raw)
        elif ext == ".pptx":
            texto = _extract_pptx(raw)
        elif ext in (".xlsx", ".xlsm"):
            texto = _extract_xlsx(raw)
        elif ext in (".zip",) or mime in ("application/zip", "application/x-zip-compressed"):
            texto = _extract_zip(raw)
        elif ext in TEXT_EXTENSIONS or mime.startswith("text/") or mime == "application/json":
            texto = _extract_plain_text(raw)
        else:
            # Tipos sem extrator dedicado (áudio, vídeo, binários, .doc/.ppt/.xls
            # antigos, ficheiros comprimidos que não .zip, etc.) — não são
            # processados nesta versão, apenas referenciados pelo nome.
            return None
    except Exception as exc:  # noqa: BLE001
        return f"[Não foi possível ler o conteúdo de '{filename}': {exc}]"

    texto = (texto or "").strip()
    if not texto:
        return None
    return _truncate(texto)